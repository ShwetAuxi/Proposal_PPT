import os

from flask import Flask, render_template, request, session, redirect, url_for, send_file
import openai
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.util import Inches

from forms import ResultsForm, GenerationForm

# insert your openai key here. Don't upload your key to GitHub or it'll get revoked!
openai.api_key = "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"

app = Flask(__name__)
bootstrap = Boostrap5(app)
app.config.update(SECRET_KEY=os.urandom(12))

@app.route("/")
def hello_world():
    # return render_template("index.html", title="Hello")
    return redirect(url_for("interactiveUI"))


CHAT_MODEL = "gpt-3.5-turbo"

MODEL = "text-davinci-003"

TEMP = 0.5

MAX_TOKENS = 100


# @app.route("/")
# def hello_world():
#     # return render_template("index.html", title="Hello")
#     return redirect(url_for("formPage"))


def generate_chat(docType, clientType, firmType, subject):
    ...


def generate(docType, clientType, firmType, subject=None):
    try:
        prompt = "Document Type: " + docType + "Client: " + clientType + "Firm: " + firmType + '\n'

        promptPres = prompt + "Given this context, generate a title for my presentation. No quotation marks:"

        promptExecSummary = prompt + "Given this context, generate an executive summary for my presentation:"

        promptAgenda = prompt + "Given this context, generate an agenda outline for my presentation:"

        session["titlePres"] = \
            openai.ChatCompletion.create(prompt=promptPres, model=MODEL, temperature=TEMP, max_tokens=MAX_TOKENS)[
                'choices'][
                0]['text']
        session["execSummary"] = \
            openai.ChatCompletion.create(prompt=promptExecSummary, model=MODEL, temperature=TEMP, max_tokens=MAX_TOKENS)[
                'choices'][0]['text']
        session["agenda"] = \
            openai.ChatCompletion.create(prompt=promptAgenda, model=MODEL, temperature=TEMP, max_tokens=MAX_TOKENS)[
                'choices'][
                0]['text']
        return redirect(url_for("results_page"))
    except openai.error.AuthenticationError as e:
        flash("OPENAI AUTH KEY EXPIRED")
        return redirect(url_for('hello_world'))
    # with the provided text, we can now create a slides with the title presentation, executive summary, and agenda

@app.route("/results_page", methods=["GET", "POST"])
def results_page():
    if request.method == "GET":
        # there's a better way to do this... but that's for later.
        if "titlePres" in session and "execSummary" in session and "agenda" in session and "pov" in session and "maras" in session:
            titlePres = session["titlePres"]
            execSummary = session["execSummary"]
            agenda = session["agenda"]
            pov = session["pov"]
            maras = session["maras"]
            return render_template("results_page.html", titlePres=titlePres, execSummary=execSummary, agenda=agenda, pov=pov, maras=maras)
        else:
            #TODO: how did we get here?
            return "SESSION DATA INVALID"
    elif request.method == "POST":
        #TODO: handle request from form
        userSatisfaction = request.form["userSatisfaction"]
        if userSatisfaction.lower() == "y":

            #return redirect(url_for("user_pref"))
            return redirect(url_for("user_pref"))
        else:
            return redirect(url_for("interactiveUI"))
        
@app.route("/user_pref", methods=["GET", "POST"])
def user_pref():
    if request.method == "POST":
        section_1 = request.form["section_1"]
        section_2 = request.form["section_2"]
        section_3 = request.form["section_3"]
        section_4 = request.form["section_4"]
        section_5 = request.form["section_5"]
        
        if section_1=="A":
            section_1_layout_index = 0  # replace with the index of the layout for option A in Title
            #elif section_1_choice == "B":
        elif section_1 == "B":
            section_1_layout_index = 1  # replace with the index of the layout for option B in Title
        else:
            print("Invalid choice. Exiting.")
            exit()
        if section_2 == "C":
            section_2_layout_index = 5  # replace with the index of the layout for option A in Agenda
        elif section_2 == "D":
            section_2_layout_index = 6  # replace with the index of the layout for option B in Agenda
        else:
            print("Invalid choice. Exiting.")
            exit()
        if section_3 == "E":
            section_3_layout_index = 5  # replace with the index of the layout for option A in Agenda
        elif section_3 == "F":
            section_3_layout_index = 6  # replace with the index of the layout for option B in Agenda
        else:
            print("Invalid choice. Exiting.")
            exit()
        if section_4 == "G":
            section_4_layout_index = 5  # replace with the index of the layout for option A in Agenda
        elif section_4 == "H":
            section_4_layout_index = 6  # replace with the index of the layout for option B in Agenda
        else:
            print("Invalid choice. Exiting.")
            exit()
        if section_5 == "I":
            section_5_layout_index = 5  # replace with the index of the layout for option A in Agenda
        elif section_5 == "J":
            section_5_layout_index = 6  # replace with the index of the layout for option B in Agenda
        else:
            print("Invalid choice. Exiting.")
            exit()               
        # Add the selected slides to a new presentation
        prs1 = Presentation("Auxi_PPT_TEMP.pptx")
        slide_layouts = prs1.slide_layouts
       
        print("Slides in slide layout")
        for i, slide in enumerate(slide_layouts):
            print(i, slide.name)
        print("Placeholders in slide layout 2")    
        for i, shape in enumerate(slide_layouts[section_2_layout_index].placeholders):
            print(i, shape.name)
        # Add a new blank text box to the slide    
        new_prs=Presentation()
        slide_1 = new_prs.slides.add_slide(slide_layouts[section_1_layout_index])
        slide_2 = new_prs.slides.add_slide(slide_layouts[section_2_layout_index])
        slide_3 = new_prs.slides.add_slide(slide_layouts[section_3_layout_index])
        slide_4 = new_prs.slides.add_slide(slide_layouts[section_4_layout_index])
        slide_5 = new_prs.slides.add_slide(slide_layouts[section_5_layout_index])
        # Set the title and content for each slide
        slide_1_content = slide_1.placeholders[1]
        slide_1_content.text = session["titlePres"]

        slide_2.shapes.title.text = "Agenda"
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(6)
        height = Inches(4)
        textbox =slide_2.shapes.add_textbox(left, top, width, height)
        textbox.text=session["agenda"]
        
        #slide_3_content = slide_3.placeholders[1]
        #slide_3_content.text = session["execSummary"]

        slide_3.shapes.title.text = "Executive Summary"
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(6)
        height = Inches(4)
        textbox1 =slide_3.shapes.add_textbox(left, top, width, height)
        #textbox1.text=session["execSummary"]

        # Get the text frame of the textbox
        text_frame = textbox1.text_frame

        # Set the text in the textbox to the contents of session["maras"]
        text_frame.text = session["execSummary"]

        #slide_4_content = slide_4.placeholders[1]
        #slide_4_content.text = session["pov"]

        slide_4.shapes.title.text = "Point of View"
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(6)
        height = Inches(4)
        textbox2 =slide_4.shapes.add_textbox(left, top, width, height)
        textbox2.text=session["pov"]

        ##slide_5_content = slide_5.placeholders[1]
        #slide_5_content.text = session["maras"]

        slide_5.shapes.title.text = "Market Assessment"
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(6)
        height = Inches(4)
        textbox3 =slide_5.shapes.add_textbox(left, top, width, height)
        textbox3.text=session["maras"]

        new_prs.save("custom_presentation.pptx")
        return "Presentation created successfully!" and send_file('custom_presentation.pptx')   
    else:
        # handle GET requests here
        return render_template("user_pref.html")    

@app.route("/form", methods=["GET", "POST"])
def formPage():
    form = GenerationForm(request.form)
    if request.method == "POST" and form.validate():
        if form.model.data == 'davinci':
            return generate(form.doc_type.data, form.client_type.data, form.firm_type.data, form.subject.data)
        elif form.model.data == 'gpt3.5':
            return generate_chat(form.doc_type.data, form.client_type.data, form.firm_type.data, form.subject.data)
        else:
            flash('INVALID CHOICE')
            return redirect(url_for('formPage'))
    return render_template('presentation_form_page.html', form=form)


@app.route("/form_old", methods=["GET", "POST"])
def interactiveUI():
    if request.method == "GET":
        # we should have the options in the form be sent to the template instead of hard coded, but we can do that later
        return render_template("presentation_form.html")
    elif request.method == "POST":
        docType = request.form["docType"]
        clientType = request.form["clientType"]
        firmType = request.form["firmType"]

        if docType == "":
            docType = request.form["docTypeOther"].lower()

        return generate(docType, clientType, firmType)
    else:
        return "INVALID REQUEST"

@app.route("/shutdown")
def shutdown():
    # this doesn't work...
    print('shutdown')
    raise RuntimeError("shutdown")

if __name__ == "__main__":
    port = int(os.environ.get('PORT', 5000))
    app.run(host=('0.0.0.0'), port=port)
